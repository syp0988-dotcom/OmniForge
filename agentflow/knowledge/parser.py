"""Document parsing pipeline supporting PDF, DOCX, TXT, Markdown, HTML,
Excel, PowerPoint, CSV, EPUB, and source code files.

Upgraded to use structure-aware chunking (see ``agentflow.knowledge.chunking``).
"""

from __future__ import annotations

import csv
import io
import re
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import IO

from agentflow.knowledge.chunking import chunk_document


def parse_document(
    file_path: str | Path,
    file_type: str | None = None,
    chunk_size: int = 500,
    chunk_overlap: int = 50,
) -> list[str]:
    """Parse a document and return a list of text chunks.

    Uses structure-aware chunking (heading-based for Markdown, definition-
    boundary for code, paragraph-based for others).

    Args:
        file_path: Path to the document file.
        file_type: Optional override (e.g. ``"pdf"``, ``"docx"``, ``"txt"``, ``"md"``).
                   Auto-detected from extension if not provided.
        chunk_size: Target character count per chunk.
        chunk_overlap: Character overlap between consecutive chunks.

    Returns:
        A list of text chunks from the document.
    """
    path = Path(file_path)
    if file_type is None:
        file_type = path.suffix.lstrip(".").lower()

    raw_text = _read_raw(path, file_type)
    return chunk_document(raw_text, file_type, chunk_size=chunk_size, overlap=chunk_overlap)


def _read_raw(path: Path, file_type: str) -> str:
    """Read the raw text content from a file based on its type."""
    ft = file_type.lower()
    if ft == "pdf":
        return _read_pdf(path)
    elif ft == "docx":
        return _read_docx(path)
    elif ft in ("md", "markdown"):
        return _read_markdown(path)
    elif ft in ("html", "htm"):
        return _read_html(path)
    elif ft in ("xlsx", "xls"):
        return _read_excel(path)
    elif ft == "pptx":
        return _read_pptx(path)
    elif ft == "csv":
        return _read_csv(path)
    elif ft == "epub":
        return _read_epub(path)
    else:
        return _read_text(path)


def _read_pdf(path: Path) -> str:
    """Extract text from a PDF file using pypdf."""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages: list[str] = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        return _fallback_read(path)


def _read_docx(path: Path) -> str:
    """Extract text from a Word document using python-docx."""
    try:
        from docx import Document

        doc = Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras)
    except ImportError:
        return _fallback_read(path)


def _read_markdown(path: Path) -> str:
    """Read a Markdown file, stripping frontmatter."""
    text = path.read_text(encoding="utf-8")
    return _strip_frontmatter(text)


def _read_text(path: Path) -> str:
    """Read a plain text file with encoding auto-detection."""
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _read_html(path: Path) -> str:
    """Extract readable text from an HTML file using BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return _fallback_read(path)
    soup = BeautifulSoup(path.read_text("utf-8", errors="replace"), "html.parser")
    # Remove non-content elements
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    lines: list[str] = []
    for el in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "pre", "blockquote"]):
        text = el.get_text(strip=True)
        if text:
            tag = el.name
            if tag.startswith("h"):
                level = tag[1]
                lines.append(f"{'#' * int(level)} {text}")
            elif tag in ("td", "th"):
                lines.append(f"[{tag}] {text}")
            elif tag == "li":
                lines.append(f"- {text}")
            elif tag == "pre":
                lines.append(f"```\n{text}\n```")
            elif tag == "blockquote":
                lines.append(f"> {text}")
            else:
                lines.append(text)
    return "\n\n".join(lines)


def _read_excel(path: Path) -> str:
    """Extract text from an Excel (.xlsx) file using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        return _fallback_read(path)
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = [f"## Sheet: {sheet_name}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("| " + " | ".join(cells) + " |")
        parts.append("\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    """Extract text from a PowerPoint (.pptx) file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return _fallback_read(path)
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append("| " + " | ".join(cells) + " |")
        parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def _read_csv(path: Path) -> str:
    """Read a CSV file and format as a markdown-style table."""
    try:
        text = path.read_text("utf-8")
    except (UnicodeDecodeError, UnicodeError):
        text = path.read_text("gbk", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""
    lines: list[str] = []
    for row in rows:
        cells = [c.strip() for c in row]
        if any(c for c in cells):
            lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _read_epub(path: Path) -> str:
    """Extract text from an EPUB file."""
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        return _fallback_read(path)
    try:
        book = epub.read_epub(str(path))
    except Exception:
        return _fallback_read(path)
    parts: list[str] = []
    from ebooklib.epub import LINK, LINK_ITEM
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            content = item.get_content()
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                continue
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text, "html.parser")
                body = soup.find("body")
                if body:
                    lines: list[str] = []
                    for el in body.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li"]):
                        t = el.get_text(strip=True)
                        if t:
                            lines.append(t)
                    if lines:
                        parts.append("\n\n".join(lines))
            except ImportError:
                import re
                clean = re.sub(r"<[^>]+>", " ", text)
                clean = re.sub(r"\s+", " ", clean).strip()
                if clean:
                    parts.append(clean)
    return "\n\n---\n\n".join(parts)


def _read_zip(path: Path) -> list[tuple[str, str]]:
    """Extract a ZIP archive and return (filename, raw_text) pairs.

    Only processes files with supported extensions.  Returns an empty
    list if none of the archive contents are parseable.
    """
    supported_exts = {
        ".pdf", ".docx", ".doc", ".txt", ".md", ".markdown",
        ".html", ".htm", ".xlsx", ".xls", ".pptx", ".csv", ".epub",
    }
    code_exts = {".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp"}
    supported_exts.update(code_exts)
    results: list[tuple[str, str]] = []
    with zipfile.ZipFile(path) as zf:
        for info in zf.infolist():
            if info.filename.startswith("__MACOSX/") or info.filename.startswith("."):
                continue
            ext = Path(info.filename).suffix.lower()
            if ext not in supported_exts:
                continue
            try:
                raw = zf.read(info.filename)
                ftype = ext.lstrip(".")
                text = _read_raw_from_bytes(raw, ftype)
                if text.strip():
                    results.append((Path(info.filename).name, text))
            except Exception:
                continue
    return results


def _read_raw_from_bytes(data: bytes, file_type: str) -> str:
    """Read raw text from bytes, dispatching by file_type.

    For simple text-based types, decode the bytes directly.
    For binary types (pdf, docx, xlsx, pptx), write to a temp file and parse.
    """
    text_types = {"txt", "md", "markdown", "html", "htm", "csv"}
    code_types = {"py", "js", "ts", "jsx", "tsx", "java", "go", "rs", "c", "cpp", "h", "hpp"}
    if file_type in text_types | code_types:
        for enc in ("utf-8", "gbk", "latin-1"):
            try:
                return data.decode(enc)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return data.decode("utf-8", errors="replace")
    # Binary types: write to temp file
    ext = "." + file_type
    with TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir) / ("tmp" + ext)
        tmp.write_bytes(data)
        return _read_raw(tmp, file_type)


def _fallback_read(path: Path) -> str:
    """Fallback: read file as plain text when a parser library is missing."""
    return _read_text(path)


def _strip_frontmatter(text: str) -> str:
    """Strip YAML/TOML frontmatter delimited by --- or +++."""
    if re.match(r"^(---|\+\+\+)\s*$", text.splitlines()[0] if text else ""):
        parts = re.split(r"^(---|\+\+\+)\s*$", text, maxsplit=2, flags=re.MULTILINE)
        if len(parts) >= 4:
            return parts[3].strip()
    return text


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Split text into overlapping chunks by paragraph boundaries.

    .. deprecated::
       Use ``agentflow.knowledge.chunking.chunk_by_paragraph`` instead.
    """
    from agentflow.knowledge.chunking import chunk_by_paragraph
    return chunk_by_paragraph(text, chunk_size, overlap)
